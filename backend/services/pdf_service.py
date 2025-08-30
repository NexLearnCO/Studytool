import PyPDF2
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except Exception:
    fitz = None
    PYMUPDF_AVAILABLE = False
from io import BytesIO
import base64
import mimetypes
import os
import hashlib
import re
from .image_storage_service import ImageStorageService

# Import for DOCX/PPTX support
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not installed, DOCX files will not be supported")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed, PPTX files will not be supported")

class PDFService:
    @staticmethod
    def extract_text(file):
        """Extract text from PDF file"""
        try:
            # Read PDF file
            pdf_reader = PyPDF2.PdfReader(BytesIO(file.read()))
            
            # Extract text from all pages
            text = ""
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"
            
            # Clean up text
            text = text.strip()
            
            if not text:
                raise ValueError("No text found in PDF")
            
            return text
            
        except Exception as e:
            raise Exception(f"Failed to extract PDF text: {str(e)}")
    
    @staticmethod
    def extract_text_from_file(file_info):
        """Extract text from file info (base64 encoded data) - supports multiple formats"""
        try:
            # Get file info
            file_name = file_info.get('name', 'unknown')
            file_data = file_info.get('data', '')
            
            # Get file extension
            _, ext = os.path.splitext(file_name.lower())
            ext = ext.lstrip('.')  # Remove the dot
            
            # Remove data URL prefix if present
            if file_data.startswith('data:'):
                file_data = file_data.split(',')[1]
            
            # Decode base64 data
            file_bytes = base64.b64decode(file_data)
            file_buffer = BytesIO(file_bytes)
            
            # Handle different file types
            if ext == 'pdf':
                return PDFService._extract_pdf_text(file_buffer)
            elif ext in ['docx', 'doc'] and DOCX_AVAILABLE:
                return PDFService._extract_docx_text(file_buffer)
            elif ext in ['pptx', 'ppt'] and PPTX_AVAILABLE:
                return PDFService._extract_pptx_text(file_buffer)
            elif ext in ['txt', 'md', 'markdown']:
                return PDFService._extract_text_file(file_buffer)
            else:
                raise ValueError(f"Unsupported file type: .{ext}")
            
        except Exception as e:
            raise Exception(f"Failed to extract text from {file_name}: {str(e)}")
    
    @staticmethod
    def _extract_pdf_text(file_buffer):
        """Extract text from PDF buffer"""
        pdf_reader = PyPDF2.PdfReader(file_buffer)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text() + "\n\n"
        
        text = text.strip()
        if not text:
            raise ValueError("No text found in PDF")
        return text

    @staticmethod
    def extract_chunks_with_metadata(file_info):
        """Use PyMuPDF to extract per-page chunks with bbox and image refs.
        Returns: list of chunks [{id, kind, text|image, doc_id, page, bbox, url?}]
        """
        if not PYMUPDF_AVAILABLE:
            # Graceful degrade: return empty list; caller may fallback
            raise ImportError("PyMuPDF (fitz) not installed. Set PYMUPDF_AVAILABLE or install PyMuPDF.")
        file_name = file_info.get('name', 'unknown')
        file_data = file_info.get('data', '')
        if file_data.startswith('data:'):
            file_data = file_data.split(',')[1]
        file_bytes = base64.b64decode(file_data)
        doc = fitz.open(stream=file_bytes, filetype='pdf')
        chunks = []
        doc_id = f'file:{file_name}'
        # Derive a filesystem-safe folder name for assets
        safe_doc_folder = re.sub(r'[^A-Za-z0-9._-]+', '_', os.path.splitext(os.path.basename(file_name))[0] or 'doc')
        assets_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'static', 'assets', safe_doc_folder))
        try:
            os.makedirs(assets_root, exist_ok=True)
        except Exception:
            # If we fail to create directory, we will still return inline placeholders
            assets_root = None
        chunk_idx = 0
        for page_index, page in enumerate(doc, 1):
            # Text blocks
            blocks = page.get_text("blocks") or []
            for b in blocks:
                x0, y0, x1, y1, text, *_ = b
                t = (text or '').strip()
                if not t:
                    continue
                chunks.append({
                    'id': f'{doc_id}-p{page_index}-c{chunk_idx}',
                    'kind': 'text',
                    'text': t,
                    'doc_id': doc_id,
                    'page': page_index,
                    'bbox': [x0, y0, x1, y1]
                })
                chunk_idx += 1
            # Images
            images = page.get_images(full=True) or []
            for img in images:
                xref = img[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    # Convert to RGB if needed (e.g., CMYK) for PNG export
                    if pix.n > 4:
                        pix_converted = fitz.Pixmap(fitz.csRGB, pix)
                        pix = pix_converted
                    image_bytes = pix.tobytes("png")
                    # Compute stable asset id by hash (dedupe)
                    asset_hash = hashlib.sha256(image_bytes).hexdigest()
                    asset_id = asset_hash[:32]
                    file_name_png = f"{asset_id}.png"
                    url = f"/static/assets/{safe_doc_folder}/{file_name_png}"
                    
                    # Extract context around image position
                    context_text = ""
                    try:
                        # Get image bbox for context extraction
                        img_rect = None
                        for block in page.get_images():
                            if block[0] == xref:
                                img_instances = page.get_image_instances(xref)
                                if img_instances:
                                    img_rect = img_instances[0]
                                    break
                        
                        # Extract text near the image
                        if img_rect:
                            # Expand search area around image
                            expand = 50  # pixels
                            search_rect = fitz.Rect(
                                max(0, img_rect.x0 - expand),
                                max(0, img_rect.y0 - expand),
                                min(page.rect.width, img_rect.x1 + expand),
                                min(page.rect.height, img_rect.y1 + expand)
                            )
                            nearby_text = page.get_textbox(search_rect)
                            if nearby_text:
                                context_text = nearby_text.strip()[:200]  # Limit context length
                    except Exception:
                        pass  # Continue without context if extraction fails
                    
                    # Persist if we can
                    if assets_root:
                        try:
                            file_path = os.path.join(assets_root, file_name_png)
                            if not os.path.exists(file_path):
                                with open(file_path, 'wb') as f:
                                    f.write(image_bytes)
                        except Exception:
                            # Fallback to inline placeholder if write fails
                            url = f'inline:image:{doc_id}:{page_index}:{xref}'
                    
                    # Store local file path for vision analysis
                    local_path = None
                    if assets_root:
                        local_path = os.path.join(assets_root, file_name_png)
                    
                    chunks.append({
                        'id': f'{doc_id}-p{page_index}-img{xref}',
                        'kind': 'image',
                        'doc_id': doc_id,
                        'page': page_index,
                        'bbox': None,
                        'asset_id': asset_id,
                        'url': url,
                        'local_path': local_path,  # Add local path for vision analysis
                        'width': pix.width,
                        'height': pix.height,
                        'context': context_text
                    })
                except Exception:
                    # On failure, keep previous behavior with inline placeholder
                    chunks.append({
                        'id': f'{doc_id}-p{page_index}-img{xref}',
                        'kind': 'image',
                        'doc_id': doc_id,
                        'page': page_index,
                        'bbox': None,
                        'url': f'inline:image:{doc_id}:{page_index}:{xref}'
                    })
        return chunks
    
    @staticmethod
    def _extract_docx_text(file_buffer):
        """Extract text from DOCX buffer"""
        if not DOCX_AVAILABLE:
            raise ValueError("DOCX support not available - install python-docx")
        
        doc = Document(file_buffer)
        text_parts = []
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip())
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        text = "\n\n".join(text_parts)
        if not text:
            raise ValueError("No text found in DOCX")
        return text
    
    @staticmethod
    def _extract_pptx_text(file_buffer):
        """Extract text from PPTX buffer"""
        if not PPTX_AVAILABLE:
            raise ValueError("PPTX support not available - install python-pptx")
        
        prs = Presentation(file_buffer)
        slides_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            slide_texts.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                elif shape.has_table:
                    # Extract table content
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_texts.append(" | ".join(row_text))
            
            if len(slide_texts) > 1:  # More than just the slide header
                slides_text.append("\n".join(slide_texts))
        
        text = "\n\n".join(slides_text)
        if not text:
            raise ValueError("No text found in PPTX")
        return text
    
    @staticmethod
    def _extract_text_file(file_buffer):
        """Extract text from plain text files"""
        try:
            # Try UTF-8 first
            text = file_buffer.getvalue().decode('utf-8')
        except UnicodeDecodeError:
            # Fallback to other encodings
            try:
                file_buffer.seek(0)
                text = file_buffer.getvalue().decode('utf-8', errors='ignore')
            except:
                file_buffer.seek(0)
                text = file_buffer.getvalue().decode('latin-1', errors='ignore')
        
        text = text.strip()
        if not text:
            raise ValueError("No text found in file")
        return text